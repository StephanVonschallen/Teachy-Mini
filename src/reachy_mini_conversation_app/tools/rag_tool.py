"""RAG tool — indexes and searches PDF/PPTX lecture slides."""

import logging
import numpy as np
from pathlib import Path
from typing import Any, Dict

from reachy_mini_conversation_app.tools.core_tools import Tool, ToolDependencies

logger = logging.getLogger(__name__)

# Global in-memory store (per session)
_index: Any = None
_chunks: list[str] = []
_source_name: str = ""
_document_context: str = ""


def get_document_context() -> str:
    """Return document context for session injection."""
    return _document_context

EMBED_MODEL = "text-embedding-3-small"
EMBED_DIM = 1536
TOP_K = 4


def _embed(texts: list[str]) -> np.ndarray:
    from openai import OpenAI
    client = OpenAI()
    resp = client.embeddings.create(model=EMBED_MODEL, input=texts)
    return np.array([r.embedding for r in resp.data], dtype="float32")


def ingest_document(file_path: str) -> str:
    logger.info(f"ingest_document called with: {file_path}")
    """Parse PDF or PPTX and build FAISS index."""
    global _index, _chunks, _source_name
    import faiss

    path = Path(file_path)
    _source_name = path.name
    raw_chunks: list[str] = []

    if path.suffix.lower() == ".pdf":
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                raw_chunks.append(f"[Page {i+1}]\n{text.strip()}")
    elif path.suffix.lower() in (".pptx", ".ppt"):
        logger.info("Processing PPTX file")
        from pptx import Presentation
        prs = Presentation(str(path))
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                raw_chunks.append(f"[Slide {i+1}]\n" + "\n".join(texts))
    else:
        return f"Unsupported file type: {path.suffix}"

    logger.info(f"raw_chunks count: {len(raw_chunks)}")
    if not raw_chunks:
        return "No text content found in document."

    _chunks = raw_chunks
    embeddings = _embed(raw_chunks)
    _index = faiss.IndexFlatL2(EMBED_DIM)
    _index.add(embeddings)

    global _document_context
    preview = _chunks[:30]
    _document_context = f"Student uploaded: '{_source_name}'. Content:\n\n" + "\n\n---\n\n".join(preview)
    return f"Document '{_source_name}' loaded: {len(_chunks)} sections indexed and ready."


class RagSearch(Tool):
    """Search uploaded lecture slides for relevant content."""

    name = "rag_search"
    description = (
        "Search the student's uploaded lecture slides or documents for relevant content. "
        "Use this when the student asks about specific topics from their course materials, "
        "or when you want to ground your explanation in their actual slide content."
    )
    parameters_schema = {
        "type": "object",
        "properties": {
            "query": {
                "type": "string",
                "description": "What to search for in the lecture materials.",
            }
        },
        "required": ["query"],
    }

    async def __call__(self, deps: ToolDependencies, **kwargs: Any) -> Dict[str, Any]:
        """Search the indexed document."""
        import faiss

        query = kwargs.get("query", "")
        if not query:
            return {"error": "Query is required"}

        if _index is None or not _chunks:
            return {"result": "No document uploaded yet. Please upload your lecture slides first."}

        q_embed = _embed([query])
        distances, indices = _index.search(q_embed, TOP_K)

        results = []
        for idx in indices[0]:
            if idx < len(_chunks):
                results.append(_chunks[idx])

        if not results:
            return {"result": "No relevant content found in the uploaded document."}

        context = "\n\n---\n\n".join(results)
        return {"result": f"Relevant content from '{_source_name}':\n\n{context}"}
